from io import BytesIO
from typing import List
from pptx import Presentation


def generate_pack(workshop: dict) -> bytes:
  prs = Presentation()
  title_slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(title_slide_layout)
  slide.shapes.title.text = workshop.get('name', 'Workshop')
  slide.placeholders[1].text = f"Sponsor: {workshop.get('sponsor','N/A')}\nScope: {', '.join(workshop.get('scope', []))}"

  heatmap = prs.slides.add_slide(prs.slide_layouts[1])
  heatmap.shapes.title.text = "Heatmap"
  body = heatmap.placeholders[1].text_frame
  body.text = "Critical gaps: {}\nHigh gaps: {}".format(
    len([g for g in workshop.get('gaps', []) if g.get('severity') == 'critical']),
    len([g for g in workshop.get('gaps', []) if g.get('severity') == 'high'])
  )

  action_slide = prs.slides.add_slide(prs.slide_layouts[1])
  action_slide.shapes.title.text = "Action Plan"
  tf = action_slide.placeholders[1].text_frame
  for action in workshop.get('actions', [])[:10]:
    p = tf.add_paragraph()
    p.text = f"{action.get('title')} — {action.get('owner')} ({action.get('status')})"

  stream = BytesIO()
  prs.save(stream)
  return stream.getvalue()
