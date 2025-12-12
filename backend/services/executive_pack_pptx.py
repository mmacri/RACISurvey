from pathlib import Path
from typing import Dict

from pptx import Presentation

PPTX_DIR = Path(__file__).resolve().parent.parent / "data" / "exports"
PPTX_DIR.mkdir(parents=True, exist_ok=True)

SLIDE_TITLES = [
    "Title & workshop context",
    "Scope & participants",
    "Current-state findings",
    "Completion heatmap by domain",
    "Ownership gaps",
    "Conflict gaps",
    "Role overload view",
    "Leadership decisions needed",
    "Action plan & timeline",
    "Method & definitions",
]


def build_pptx(summary: Dict, workshop_id: int) -> Path:
    prs = Presentation()
    for title in SLIDE_TITLES:
        slide_layout = prs.slide_layouts[0] if prs.slide_layouts else prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{title}"
        if slide.placeholders:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:  # body placeholder
                    placeholder.text = f"Auto-generated summary: {summary}"
    target = PPTX_DIR / f"workshop_{workshop_id}_executive.pptx"
    prs.save(target)
    return target
